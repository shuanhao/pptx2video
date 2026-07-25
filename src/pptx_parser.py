from pathlib import Path
from typing import List, Dict, Optional

from pptx import Presentation


def _get_slide_title(slide) -> Optional[str]:
    title = None
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text_frame = shape.text_frame
            if text_frame and text_frame.text.strip():
                text = text_frame.text.strip()
                if text and text.lower() != "click to add title":
                    title = text
                    break
    return title


def _get_notes_text(slide) -> Optional[str]:
    notes_text = ""
    try:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
    except Exception:
        return None

    if text_frame is None:
        return None

    paragraphs = []
    for paragraph in text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        paragraphs.append(text)

    notes_text = "\n".join(paragraphs).strip()
    return notes_text or None


def extract_notes(pptx_path: str) -> List[Dict[str, Optional[str]]]:
    """Extract slide number, title, and notes text from a PowerPoint file.

    Args:
        pptx_path: Path to the .pptx file.

    Returns:
        A list of dictionaries with 'slide_num', 'title', and 'notes'.
    """
    path = Path(pptx_path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

    if path.suffix.lower() != ".pptx":
        raise ValueError(f"Unsupported file format: {path.suffix}")

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        raise RuntimeError(f"Failed to load PowerPoint file: {exc}") from exc

    slides_data: List[Dict[str, Optional[str]]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = _get_slide_title(slide)
        notes = _get_notes_text(slide)
        slides_data.append({
            "slide_num": idx,
            "title": title,
            "notes": notes,
        })

    return slides_data


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage: python pptx_parser.py <path-to-pptx>")
        sys.exit(1)

    try:
        slides = extract_notes(sys.argv[1])
        print(f"Parsed {len(slides)} slide(s) from {sys.argv[1]}")
        for item in slides:
            title = item["title"] or "[no title]"
            notes = item["notes"] or "[no notes]"
            print(f"- Slide {item['slide_num']}: {title}")
            print(f"  Notes: {notes}")
    except FileNotFoundError as exc:
        print(f"Error: {exc}")
        sys.exit(1)
    except ValueError as exc:
        print(f"Error: {exc}")
        sys.exit(1)
    except RuntimeError as exc:
        print(f"Error: {exc}")
        sys.exit(1)
    except Exception as exc:
        print(f"Unexpected error: {exc}")
        sys.exit(1)
