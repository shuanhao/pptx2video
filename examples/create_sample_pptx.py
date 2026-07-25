from pathlib import Path
from pptx import Presentation

path = Path(__file__).with_name('sample_test.pptx')
prs = Presentation()

slides_data = [
    {
        'title': 'Cover Page',
        'notes': [],
    },
    {
        'title': 'Page 1',
        'notes': [
            '這一段我想先講一個很簡單的概念，因為這個例子本來就是為了讓我們看見，真的有可能把簡報內容變成一段很自然的演講稿。',
            '我會用一個英文單字 example，然後配上一句 Hello World，讓它看起來更貼近你平常在做簡報時，會自然說出來的內容。',
            '',
            '接下來這一段會換很多行，包含中文描述、English phrase，還有一些比較長的敘述，這樣才比較像我們在現場講話時的節奏。',
            '例如：這裡要測試 line break、paragraph break，還有一段非常長的內容，因為講稿如果太短，聽起來就不太像真的在演講。',
            '我希望這個 parser 在讀取時，能夠正確保留換行與空白行，並且不會把中英文混在一起的內容吃掉，這樣後面才有辦法變成自然的語音。',
            '',
            '最後再補一段更完整的內容，像是 product demo、AI workflow、next step 這種混合語境的說明，這樣聽起來就更像一位講者在說明他正在做的事情。',
            '如果你在做正式簡報，這種內容就很像是講者自己準備的口語稿，會比較容易被 TTS 讀得自然。',
        ],
    },
    {
        'title': 'Page 2',
        'notes': [
            '這一頁是另一個較長的示範，內容會包含更多中文段落，並且夾雜一些英文單字，例如 API、UI、workflow、token。',
            '我們希望測試在多段落中，是否能正確保留換行、空白行，以及不同語系混排時的順序，因為這些細節會直接影響你之後的語音呈現。',
            '',
            '這裡也可以放一段更長的內容，像是我們在現場講解一個產品時，會自然說出來的話。',
            '第一行：這是一個很長的說明，目的是讓你看見如何在 notes 裡面編排長文，這樣講的時候才不會一口氣講太快。',
            '第二行：這裡又加入了英文縮寫 like MVP、beta release、demo version，這種混合語彙其實很常出現在真實簡報裡。',
            '第三行：我們也想確認是否可以同時處理中文與英文，像是 AI model、cloud service、real-time update，這些詞彙都很容易出現在現場說明裡。',
            '',
            '最後再加上一個更長的結尾，像是：',
            '這是為了驗證 notes 在多行與多段落時，是否仍然能保持原始格式，並且不會被誤刪除空白行，因為這對於後續生成語音來說真的很重要。',
            '如果你把這種內容直接拿去做 TTS，聽起來就會很像一段真正的講稿，而不是一堆機械化的文字。',
        ],
    },
    {
        'title': 'Page 3',
        'notes': [
            '這是第三頁的測試內容，特別強調長篇筆記與中英混排的效果，因為這正是很多講者在現場最常遇到的情況。',
            '例如：這裡有一個英文名詞 design system，還有中文解釋，讓整段內容更接近實際簡報使用情境，也更容易讓人理解。',
            '',
            '我們也可以再加一些較長的句子，像是：',
            'In this section, we want to verify whether the parser can preserve multiple lines and blank lines correctly, so the final speech sounds natural and not robotic.',
            '同時也要確認中文詞彙與英文單字之間的相對位置，像是 feature、summary、analysis、結論，這些其實都很常出現在正式演講裡。',
            '',
            '最後一段則是用來測試更長的排版空間，包含很多行、很多段落、還有一些中英文夾雜的內容，例如：',
            'iteration plan, review notes, follow-up action, 後續要追蹤的項目，還有一個小小的提醒：不要漏掉任何一個重要細節。',
            '這種寫法對於後續生成語音也會很有幫助，因為它更像一份可直接唸出的講稿，而不是一堆生硬的文字。',
        ],
    },
    {
        'title': 'Q&A / Thanks',
        'notes': [],
    },
]

for slide_data in slides_data:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = slide_data['title']
    slide.placeholders[1].text = f'Content for {slide_data["title"]}'

    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.clear()

    for index, paragraph_text in enumerate(slide_data['notes']):
        if index == 0:
            notes_frame.paragraphs[0].text = paragraph_text
        else:
            new_paragraph = notes_frame.add_paragraph()
            new_paragraph.text = paragraph_text

prs.save(path)
print(path.exists(), path)
