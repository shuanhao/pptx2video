import io
import json
import sys
import tempfile
import unittest
from contextlib import redirect_stderr, redirect_stdout
from pathlib import Path
from unittest import mock

from pptx import Presentation

from src.logging_config import shutdown_logging
from src.main import main


class CliEndToEndTests(unittest.TestCase):
    """End-to-end tests that invoke src.main.main() the same way the real
    CLI entry point does - patching sys.argv and letting argparse, the
    parser/TTS/subtitle orchestration, and file I/O all run for real.

    This covers a gap the rest of the test suite doesn't: every other test
    module calls individual functions directly (extract_notes,
    build_payload, generate_audio_files, ...). Nothing exercised main()
    itself, so a wiring mistake there (a wrong argument passed through, a
    step run in the wrong order, an option not actually reaching the
    function it's supposed to configure) could slip past a fully green
    test suite.

    PowerPoint automation (--insert-audio / --export-video) still can't be
    covered here since it requires real Windows + PowerPoint; those two
    flags are intentionally not exercised - see test_ppt_automation.py for
    the fake-COM-object coverage of that code path instead.
    """

    def tearDown(self):
        # main() calls setup_logging() with the shared "pptx2video" logger
        # name - the same one the real CLI uses, and the same one other
        # test modules may touch. Release its handlers after each test so
        # they don't accumulate across test methods/modules, and so the
        # next test's handler is bound to *that* test's redirected stdout
        # rather than a stale one from a previous test.
        shutdown_logging()

    def _create_pptx(self, slide_specs):
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        path = Path(temp_dir.name) / "sample.pptx"

        prs = Presentation()
        for title, notes in slide_specs:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            if title:
                try:
                    slide.shapes.title.text = title
                except AttributeError:
                    pass
            if notes:
                slide.notes_slide.notes_text_frame.text = notes

        prs.save(path)
        return path

    def _invoke(self, argv):
        """Run main() with argv patched, capturing stdout/stderr.

        Returns (stdout, stderr, exit_code). exit_code is None on a normal
        return (success path) or the SystemExit code argparse/_fail() used
        (failure path) - callers don't need two different calling
        conventions depending on whether they expect success or failure.
        """
        stdout, stderr = io.StringIO(), io.StringIO()
        exit_code = None
        with mock.patch.object(sys, "argv", ["main.py"] + argv):
            with redirect_stdout(stdout), redirect_stderr(stderr):
                try:
                    main()
                except SystemExit as exc:
                    exit_code = exc.code
        return stdout.getvalue(), stderr.getvalue(), exit_code

    def test_parses_pptx_and_writes_json_and_srt_without_audio(self):
        pptx_path = self._create_pptx([
            ("Intro", "Hello there"),
            ("Cover", None),
            ("Outro", "Thanks for watching"),
        ])

        with tempfile.TemporaryDirectory() as tmp:
            json_path = Path(tmp) / "slides.json"
            srt_path = Path(tmp) / "captions.srt"
            # Point --audio-output-dir at an empty tempdir explicitly. This
            # test asserts that *no* audio manifest means an empty .srt is
            # written - but the flag defaults to the relative "output/audio",
            # which resolves against whatever the test runner's cwd happens
            # to be. Every other test in this file isolates this same way;
            # this one didn't, so on a real checkout that's actually been
            # used to run the CLI for real (i.e. one with a genuine
            # output/audio/manifest.json left over on disk from a previous
            # run), this test would pick that manifest up, resolve real
            # (mismatched) timing data against these fixture slides, and
            # write a non-empty .srt full of zero-duration/misaligned cues -
            # exactly the failure this isolation prevents.
            audio_dir = Path(tmp) / "audio"

            _, _, exit_code = self._invoke([
                str(pptx_path),
                "--output", str(json_path),
                "--subtitles-output", str(srt_path),
                "--audio-output-dir", str(audio_dir),
                "--no-file-log",
            ])

            self.assertIsNone(exit_code)
            self.assertTrue(json_path.exists())
            payload = json.loads(json_path.read_text(encoding="utf-8"))
            self.assertEqual(payload["slide_count"], 3)
            self.assertEqual(payload["slides"][0]["title"], "Intro")
            self.assertTrue(payload["slides"][0]["has_notes"])
            self.assertFalse(payload["slides"][1]["has_notes"])
            self.assertIsNone(payload["slides"][0]["audio_file"])

            # No --generate-audio was passed and no manifest.json exists to
            # load, so there's no WordBoundary timing data to build real
            # subtitle lines from (see subtitle_pipeline.generate_srt_for_deck's
            # docstring) - an empty .srt is written rather than failing the
            # run or falling back to a lower-quality guess.
            self.assertTrue(srt_path.exists())
            self.assertEqual(srt_path.read_text(encoding="utf-8"), "")

    def test_generate_audio_flag_wires_manifest_into_json_output(self):
        pptx_path = self._create_pptx([("Intro", "Hello there")])

        fake_manifest = {
            "voice": "fake-voice",
            "rate": "-10%",
            "pitch": "+0Hz",
            "output_dir": "output/audio",
            "slides": [{"slide_num": 1, "title": "Intro", "audio_file": "slide_001.mp3"}],
        }

        with tempfile.TemporaryDirectory() as tmp:
            json_path = Path(tmp) / "slides.json"
            audio_dir = Path(tmp) / "audio"

            # Patch generate_audio_files at its src.main import site rather
            # than in src.tts, so this exercises main()'s real
            # orchestration (argument wiring, JSON assembly, progress
            # printing) without making a real network call to edge-tts -
            # the same "fake the external dependency" approach the rest of
            # the suite already uses for PowerPoint COM.
            with mock.patch(
                "src.main.generate_audio_files", return_value=fake_manifest
            ) as mock_generate:
                _, _, exit_code = self._invoke([
                    str(pptx_path),
                    "--output", str(json_path),
                    "--generate-audio",
                    "--audio-output-dir", str(audio_dir),
                    "--subtitles-output", str(Path(tmp) / "captions.srt"),
                    "--no-file-log",
                ])

            self.assertIsNone(exit_code)
            mock_generate.assert_called_once()
            _, kwargs = mock_generate.call_args
            self.assertEqual(
                kwargs["voice"],
                "Microsoft Server Speech Text to Speech Voice (zh-TW, YunJheNeural)",
            )

            payload = json.loads(json_path.read_text(encoding="utf-8"))
            self.assertEqual(payload["audio"]["voice"], "fake-voice")
            self.assertTrue(payload["slides"][0]["audio_file"].endswith("slide_001.mp3"))

    def test_slides_flag_regenerates_only_selected_slide_and_merges_manifest(self):
        # Simulates the real use case this flag exists for: a manifest.json
        # already exists from a prior full run (e.g. one that predates the
        # dropped-narration check, or where a single slide needs
        # re-generating after investigating a POSSIBLE DROPPED NARRATION
        # warning) - --slides 2 should regenerate only slide 2's audio and
        # leave slide 1's and slide 3's manifest entries from the prior run
        # untouched, rather than truncating manifest.json down to just slide 2.
        pptx_path = self._create_pptx([
            ("One", "First"),
            ("Two", "Second"),
            ("Three", "Third"),
        ])

        with tempfile.TemporaryDirectory() as tmp:
            audio_dir = Path(tmp) / "audio"
            audio_dir.mkdir()
            existing_manifest = {
                "voice": "old-voice",
                "rate": "-10%",
                "pitch": "+0Hz",
                "output_dir": str(audio_dir),
                "slides": [
                    {"slide_num": 1, "title": "One", "audio_file": "slide_001.mp3", "word_boundaries_file": None, "narration_gap_warnings": []},
                    {"slide_num": 2, "title": "Two", "audio_file": "slide_002.mp3", "word_boundaries_file": None, "narration_gap_warnings": []},
                    {"slide_num": 3, "title": "Three", "audio_file": "slide_003.mp3", "word_boundaries_file": None, "narration_gap_warnings": []},
                ],
            }
            (audio_dir / "manifest.json").write_text(json.dumps(existing_manifest), encoding="utf-8")

            regenerated_manifest = {
                "voice": "new-voice",
                "rate": "-10%",
                "pitch": "+0Hz",
                "output_dir": str(audio_dir),
                "slides": [
                    {"slide_num": 2, "title": "Two", "audio_file": "slide_002.mp3", "word_boundaries_file": "slide_002.wordboundaries.json", "narration_gap_warnings": []},
                ],
            }

            with mock.patch(
                "src.main.generate_audio_files", return_value=regenerated_manifest
            ) as mock_generate:
                _, _, exit_code = self._invoke([
                    str(pptx_path),
                    "--output", str(Path(tmp) / "slides.json"),
                    "--generate-audio",
                    "--audio-output-dir", str(audio_dir),
                    "--slides", "2",
                    "--subtitles-output", str(Path(tmp) / "captions.srt"),
                    "--no-file-log",
                ])

            self.assertIsNone(exit_code)
            # generate_audio_files should only have been asked to generate
            # slide 2, not all three.
            call_args, _ = mock_generate.call_args
            generated_slides = call_args[0]
            self.assertEqual([s["slide_num"] for s in generated_slides], [2])

            merged = json.loads((audio_dir / "manifest.json").read_text(encoding="utf-8"))
            self.assertEqual([s["slide_num"] for s in merged["slides"]], [1, 2, 3])
            self.assertEqual(merged["slides"][0]["audio_file"], "slide_001.mp3")  # untouched
            self.assertEqual(merged["slides"][1]["word_boundaries_file"], "slide_002.wordboundaries.json")  # regenerated
            self.assertEqual(merged["slides"][2]["audio_file"], "slide_003.mp3")  # untouched

    def test_slides_flag_rejects_slide_number_not_in_deck(self):
        pptx_path = self._create_pptx([("Only", "Some notes")])

        with tempfile.TemporaryDirectory() as tmp:
            _, stderr, exit_code = self._invoke([
                str(pptx_path),
                "--output", str(Path(tmp) / "slides.json"),
                "--generate-audio",
                "--audio-output-dir", str(Path(tmp) / "audio"),
                "--slides", "99",
                "--no-file-log",
            ])

            self.assertIsNotNone(exit_code)
            self.assertNotEqual(exit_code, 0)
            self.assertIn("99", stderr)

    def test_subtitles_output_alone_reuses_existing_manifest_without_generate_audio(self):
        # Regression test for a real gap found via a user question: after
        # audio was already generated in an earlier, separate
        # `--generate-audio` run (leaving a real manifest.json + per-slide
        # wordboundaries.json on disk), running *just* `--subtitles-output`
        # on its own (no `--generate-audio`, no `--export-video`) used to
        # silently write an empty .srt instead of using that already-generated
        # data - even though the exact same scenario, routed through
        # `--insert-audio` or the post-`--export-video` true-start subtitle
        # path, already correctly loaded it via `_resolve_audio_manifest()`.
        # That meant producing subtitles without re-generating audio (the
        # whole point of separating steps 7/8 to save re-paying for edge-tts)
        # only worked if you also exported a video in that same command -
        # confirmed by reproducing it directly against main(). Fixed by
        # routing the predictive (no-export-video) subtitle path through the
        # same `_resolve_audio_manifest()` fallback.
        pptx_path = self._create_pptx([("Intro", "Hello there test notes")])

        with tempfile.TemporaryDirectory() as tmp:
            audio_dir = Path(tmp) / "audio"
            audio_dir.mkdir()
            (audio_dir / "slide_001.mp3").write_bytes(b"\x00")
            word_boundaries = [
                {"text": t, "offset_seconds": i * 0.3, "duration_seconds": 0.25}
                for i, t in enumerate(["Hello", " ", "there", " ", "test", " ", "notes"])
            ]
            (audio_dir / "slide_001.wordboundaries.json").write_text(
                json.dumps(word_boundaries), encoding="utf-8"
            )
            manifest = {
                "voice": "fake-voice",
                "rate": "-10%",
                "pitch": "+0Hz",
                "output_dir": str(audio_dir),
                "slides": [{
                    "slide_num": 1,
                    "title": "Intro",
                    "audio_file": "slide_001.mp3",
                    "word_boundaries_file": "slide_001.wordboundaries.json",
                }],
            }
            (audio_dir / "manifest.json").write_text(json.dumps(manifest), encoding="utf-8")

            srt_path = Path(tmp) / "captions.srt"
            _, _, exit_code = self._invoke([
                str(pptx_path),
                "--output", str(Path(tmp) / "slides.json"),
                "--audio-output-dir", str(audio_dir),
                "--subtitles-output", str(srt_path),
                "--no-file-log",
            ])

            self.assertIsNone(exit_code)
            srt_text = srt_path.read_text(encoding="utf-8")
            self.assertNotEqual(srt_text, "")
            self.assertIn("Hello there test", srt_text)

    def test_negative_tts_max_retries_is_rejected_before_any_work_happens(self):
        # End-to-end version of the CLI-layer guard added in v0.4.1
        # (test_pptx_parser.py already checks build_parser() in isolation;
        # this confirms the rejection actually happens when main() runs,
        # i.e. before the pptx is even parsed).
        pptx_path = self._create_pptx([("Intro", "Hello there")])

        with tempfile.TemporaryDirectory() as tmp:
            json_path = Path(tmp) / "slides.json"

            _, stderr, exit_code = self._invoke([
                str(pptx_path),
                "--output", str(json_path),
                "--tts-max-retries", "-1",
                "--no-file-log",
            ])

            self.assertEqual(exit_code, 2)
            self.assertIn("--tts-max-retries", stderr)
            self.assertFalse(json_path.exists())

    def test_strict_mode_aborts_when_a_slide_has_no_notes(self):
        pptx_path = self._create_pptx([("Intro", "Hello there"), ("No notes", None)])

        with tempfile.TemporaryDirectory() as tmp:
            json_path = Path(tmp) / "slides.json"

            _, stderr, exit_code = self._invoke([
                str(pptx_path),
                "--output", str(json_path),
                "--subtitles-output", str(Path(tmp) / "captions.srt"),
                "--strict",
                "--no-file-log",
            ])

            self.assertEqual(exit_code, 2)
            self.assertIn("Strict mode", stderr)
            self.assertFalse(json_path.exists())

    def test_missing_pptx_file_exits_cleanly_with_error_message(self):
        with tempfile.TemporaryDirectory() as tmp:
            missing_path = Path(tmp) / "does-not-exist.pptx"
            json_path = Path(tmp) / "slides.json"

            _, stderr, exit_code = self._invoke([
                str(missing_path),
                "--output", str(json_path),
                "--subtitles-output", str(Path(tmp) / "captions.srt"),
                "--no-file-log",
            ])

            self.assertEqual(exit_code, 2)
            self.assertIn("error:", stderr)
            self.assertFalse(json_path.exists())

    def test_pretty_flag_prints_json_to_stdout(self):
        pptx_path = self._create_pptx([("Intro", "Hello there")])

        with tempfile.TemporaryDirectory() as tmp:
            json_path = Path(tmp) / "slides.json"
            stdout, _, exit_code = self._invoke([
                str(pptx_path),
                "--output", str(json_path),
                "--pretty",
                "--subtitles-output", str(Path(tmp) / "captions.srt"),
                "--no-file-log",
            ])

            # Console logging (e.g. "Saved JSON to ...") also writes to
            # stdout, so the captured output isn't pure JSON - assert on
            # the pretty-printed JSON appearing somewhere in it rather than
            # trying to json.loads() the whole capture.
            self.assertIsNone(exit_code)
            self.assertIn('"slide_count": 1', stdout)
            self.assertIn('"title": "Intro"', stdout)


if __name__ == "__main__":
    unittest.main()
