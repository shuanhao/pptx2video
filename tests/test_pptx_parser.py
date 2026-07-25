import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from src.main import build_parser
from src.pptx_parser import extract_notes


class PptxParserTests(unittest.TestCase):
    def _create_pptx(self, slide_specs):
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        path = Path(temp_dir.name) / "sample.pptx"

        prs = Presentation()
        for title, notes in slide_specs:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            if title:
                try:
                    slide.shapes.title.text = title
                except AttributeError:
                    pass
            if notes:
                slide.notes_slide.notes_text_frame.text = notes

        prs.save(path)
        return str(path)

    def test_extract_notes_returns_all_slides_with_titles_and_notes(self):
        pptx_path = self._create_pptx([("First slide", "First note"), ("Second slide", "Second note")])

        slides = extract_notes(pptx_path)

        self.assertEqual(len(slides), 2)
        self.assertEqual(slides[0]["title"], "First slide")
        self.assertEqual(slides[0]["notes"], "First note")
        self.assertEqual(slides[1]["title"], "Second slide")
        self.assertEqual(slides[1]["notes"], "Second note")

    def test_extract_notes_handles_missing_notes(self):
        pptx_path = self._create_pptx([("Only title", "")])

        slides = extract_notes(pptx_path)

        self.assertEqual(len(slides), 1)
        self.assertEqual(slides[0]["title"], "Only title")
        self.assertIsNone(slides[0]["notes"])

    def test_extract_notes_preserves_blank_lines_between_paragraphs(self):
        temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(temp_dir.cleanup)
        path = Path(temp_dir.name) / "blank-lines.pptx"

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        notes_frame = slide.notes_slide.notes_text_frame
        first_paragraph = notes_frame.paragraphs[0]
        first_paragraph.text = "First paragraph"
        empty_paragraph = notes_frame.add_paragraph()
        empty_paragraph.text = ""
        third_paragraph = notes_frame.add_paragraph()
        third_paragraph.text = "Third paragraph"
        prs.save(path)

        slides = extract_notes(str(path))

        self.assertEqual(slides[0]["notes"], "First paragraph\n\nThird paragraph")

    def test_build_parser_supports_verbose_and_strict_flags(self):
        parser = build_parser()
        args = parser.parse_args(["demo.pptx", "--verbose", "--strict", "--output", "out.json"])

        self.assertEqual(args.pptx_path, "demo.pptx")
        self.assertTrue(args.verbose)
        self.assertTrue(args.strict)
        self.assertEqual(args.output, "out.json")


if __name__ == "__main__":
    unittest.main()
